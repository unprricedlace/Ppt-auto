import copy
from pptx import Presentation
from pptx.enum.text import MSO_TEXT_EFFECT
from pptx.dml.color import RGBColor

def extract_text_from_shape(shape):
    """
    Extract text from various types of shapes, including text boxes, placeholders, and table cells
    """
    if shape.has_text_frame:
        return shape.text_frame.text
    elif shape.has_table:
        # Extract text from table cells
        table_text = []
        for row in shape.table.rows:
            row_text = []
            for cell in row.cells:
                row_text.append(cell.text)
            table_text.append(row_text)
        return table_text
    return ""

def compare_presentations(ppt1_path, ppt2_path, output_path):
    """
    Compare two PowerPoint presentations and create a new presentation
    with unique content highlighted in red
    """
    # Load both presentations
    prs1 = Presentation(ppt1_path)
    prs2 = Presentation(ppt2_path)
    
    # Create a new presentation as a copy of the second presentation
    new_prs = Presentation(ppt2_path)
    
    # Ensure we have the same number of slides
    min_slides = min(len(prs1.slides), len(prs2.slides))
    
    # Compare slides
    for slide_idx in range(min_slides):
        slide1 = prs1.slides[slide_idx]
        slide2 = prs2.slides[slide_idx]
        new_slide = new_prs.slides[slide_idx]
        
        # Extract text from first slide
        slide1_texts = {}
        for shape1 in slide1.shapes:
            text1 = extract_text_from_shape(shape1)
            if text1:
                slide1_texts[id(shape1)] = text1
        
        # Compare with second slide
        for shape2_idx, shape2 in enumerate(slide2.shapes):
            text2 = extract_text_from_shape(shape2)
            
            # If the shape exists in slide1 and text matches, skip
            match_found = False
            for shape1_id, text1 in slide1_texts.items():
                if text1 == text2:
                    match_found = True
                    break
            
            # If no match found, highlight the text in red
            if not match_found and text2:
                # For text frames
                if shape2.has_text_frame:
                    for paragraph in shape2.text_frame.paragraphs:
                        paragraph.font.color.rgb = RGBColor(255, 0, 0)  # Red color
                
                # For tables
                elif shape2.has_table:
                    for row in shape2.table.rows:
                        for cell in row.cells:
                            if cell.text:
                                for paragraph in cell.text_frame.paragraphs:
                                    paragraph.font.color.rgb = RGBColor(255, 0, 0)  # Red color
    
    # Save the new presentation
    new_prs.save(output_path)
    print(f"Comparison complete. New presentation saved to {output_path}")

def main():
    # Paths to the input and output presentations
    ppt1_path = 'presentation1.pptx'  # Replace with your first presentation path
    ppt2_path = 'presentation2.pptx'  # Replace with your second presentation path
    output_path = 'compared_presentation.pptx'  # Output presentation path
    
    # Run the comparison
    compare_presentations(ppt1_path, ppt2_path, output_path)

if __name__ == '__main__':
    main()

# Note: This script requires python-pptx library
# Install it using: pip install python-pptx
